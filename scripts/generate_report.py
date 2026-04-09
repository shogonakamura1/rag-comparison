"""RAG vs NotebookLM 比較実験レポート生成スクリプト

参考スライド（SDD.pptx）の色とフォントに合わせる:
- メインカラー: #2D2D8A (濃紺), #009999 (ティール), #333399 (紫青)
- アクセント: #BBE0E3 (薄水色), #DAEDEF (極薄水色), #99CC00 (黄緑)
- フォント: Meiryo UI
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

# テーマカラー
NAVY = RGBColor(0x2D, 0x2D, 0x8A)
TEAL = RGBColor(0x00, 0x99, 0x99)
PURPLE_BLUE = RGBColor(0x33, 0x33, 0x99)
LIGHT_CYAN = RGBColor(0xBB, 0xE0, 0xE3)
PALE_CYAN = RGBColor(0xDA, 0xED, 0xEF)
LIME = RGBColor(0x99, 0xCC, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x80, 0x80, 0x80)
DARK_GRAY = RGBColor(0x40, 0x40, 0x40)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
RED = RGBColor(0xCC, 0x33, 0x33)
GREEN = RGBColor(0x33, 0x99, 0x33)
ORANGE = RGBColor(0xE6, 0x80, 0x00)

FONT_JP = "Meiryo UI"

# プレゼンテーション初期化
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
SW = prs.slide_width
SH = prs.slide_height

BLANK_LAYOUT = prs.slide_layouts[6]


def set_font(run, size, bold=False, color=BLACK, font=FONT_JP):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, left, top, width, height, text, size=14, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_JP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size, bold, color, font)
    return box


def add_rect(slide, left, top, width, height, fill=LIGHT_CYAN, line=None, line_width=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        if line_width is not None:
            rect.line.width = Pt(line_width)
    return rect


def add_round_rect(slide, left, top, width, height, fill=LIGHT_CYAN, line=None, line_width=None, radius=0.1):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        if line_width is not None:
            rect.line.width = Pt(line_width)
    # adjustments
    rect.adjustments[0] = radius
    return rect


def add_arrow(slide, left, top, width, height, fill=NAVY):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = fill
    arrow.line.fill.background()
    return arrow


def add_line(slide, x1, y1, x2, y2, color=GRAY, weight=1.0):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_title_bar(slide, title_text):
    """各スライドのタイトル帯（左端の縦バー＋タイトル）"""
    bar = add_rect(slide, Inches(0.4), Inches(0.35), Inches(0.08), Inches(0.45), fill=TEAL)
    add_text(slide, Inches(0.6), Inches(0.32), Inches(9.0), Inches(0.5),
             title_text, size=22, bold=True, color=NAVY)
    add_line(slide, Inches(0.4), Inches(0.92), Inches(9.6), Inches(0.92), color=TEAL, weight=1.5)


def add_footer(slide, page_num, total):
    add_text(slide, Inches(0.4), Inches(5.35), Inches(5.0), Inches(0.25),
             "RAG vs NotebookLM 比較実験レポート", size=9, color=GRAY)
    add_text(slide, Inches(8.8), Inches(5.35), Inches(0.8), Inches(0.25),
             f"{page_num} / {total}", size=9, color=GRAY, align=PP_ALIGN.RIGHT)


# =====================================================
# 全スライド数を後で更新するためのリスト
# =====================================================
slides_def = []  # (build_func) のリスト

# =====================================================
# Slide 1: タイトル
# =====================================================
def slide_title(slide, page, total):
    # 上下の帯
    add_rect(slide, 0, 0, SW, Inches(0.6), fill=NAVY)
    add_rect(slide, 0, SH - Inches(0.4), SW, Inches(0.4), fill=NAVY)

    # 左の縦アクセント
    add_rect(slide, Inches(0.5), Inches(1.2), Inches(0.15), Inches(3.0), fill=TEAL)

    # メインタイトル
    add_text(slide, Inches(0.8), Inches(1.3), Inches(8.5), Inches(1.0),
             "RAGシステムとGoogle NotebookLM", size=32, bold=True, color=NAVY)
    add_text(slide, Inches(0.8), Inches(2.2), Inches(8.5), Inches(1.0),
             "の比較実験", size=32, bold=True, color=NAVY)

    # サブタイトル
    add_text(slide, Inches(0.8), Inches(3.3), Inches(8.5), Inches(0.5),
             "Claude 公式ドキュメントを対象とした事実抽出評価", size=16, color=DARK_GRAY)

    # 学生情報
    add_text(slide, Inches(0.8), Inches(4.4), Inches(8.5), Inches(0.4),
             "電子情報工学科　中村　翔吾", size=14, color=DARK_GRAY)
    add_text(slide, Inches(0.8), Inches(4.7), Inches(8.5), Inches(0.4),
             "2026年4月", size=12, color=GRAY)


# =====================================================
# Slide 2: 研究の目的
# =====================================================
def slide_purpose(slide, page, total):
    add_title_bar(slide, "1. 研究の目的")

    add_text(slide, Inches(0.6), Inches(1.1), Inches(9.0), Inches(0.4),
             "課題", size=14, bold=True, color=TEAL)
    add_round_rect(slide, Inches(0.6), Inches(1.45), Inches(9.0), Inches(0.65), fill=PALE_CYAN)
    add_text(slide, Inches(0.85), Inches(1.55), Inches(8.6), Inches(0.5),
             "「自分が作ったRAGシステムとGoogle NotebookLMの違いを観察する」",
             size=14, color=NAVY)

    add_text(slide, Inches(0.6), Inches(2.3), Inches(9.0), Inches(0.4),
             "実施項目", size=14, bold=True, color=TEAL)
    items = [
        ("①", "自作RAGシステムの構築（Python・ChromaDB・Gemini API）"),
        ("②", "ソース文書（10ファイル）と質問セットの作成"),
        ("③", "RAGとNotebookLMの両方で同じ質問を実行・回答を比較"),
        ("④", "RAGの弱点を分析し、改善実装で性能向上を試行"),
    ]
    y = 2.65
    for num, text in items:
        add_text(slide, Inches(0.85), Inches(y), Inches(0.4), Inches(0.4),
                 num, size=14, bold=True, color=TEAL)
        add_text(slide, Inches(1.2), Inches(y), Inches(8.4), Inches(0.4),
                 text, size=13, color=BLACK)
        y += 0.45

    add_text(slide, Inches(0.6), Inches(4.65), Inches(9.0), Inches(0.4),
             "ゴール", size=14, bold=True, color=TEAL)
    add_text(slide, Inches(0.6), Inches(4.9), Inches(9.0), Inches(0.4),
             "両システムの構造的な違いを実験により定量的に明らかにする",
             size=13, color=BLACK)


# =====================================================
# Slide 3: 実験用RAGシステムのアーキテクチャ
# =====================================================
def slide_architecture(slide, page, total):
    add_title_bar(slide, "2. 自作RAGシステムのアーキテクチャ")

    # フェーズ1: インジェスト
    add_text(slide, Inches(0.6), Inches(1.1), Inches(9.0), Inches(0.3),
             "フェーズ1: ドキュメント取り込み（インジェスト）", size=12, bold=True, color=TEAL)

    boxes_top = [
        ("URL", "Claude公式\nドキュメント"),
        ("WebFetch", "本文抽出"),
        ("チャンク分割", "500文字\n+ オーバーラップ100"),
        ("Embedding", "multilingual-\ne5-large"),
        ("ChromaDB", "ローカル\n保存"),
    ]
    box_w = 1.6
    box_h = 0.85
    gap = 0.25
    total_w = box_w * len(boxes_top) + gap * (len(boxes_top) - 1)
    start_x = (10 - total_w) / 2
    y = 1.45
    for i, (label, desc) in enumerate(boxes_top):
        x = start_x + i * (box_w + gap)
        add_round_rect(slide, Inches(x), Inches(y), Inches(box_w), Inches(box_h),
                       fill=LIGHT_CYAN, line=NAVY, line_width=0.75)
        add_text(slide, Inches(x), Inches(y + 0.05), Inches(box_w), Inches(0.3),
                 label, size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(x), Inches(y + 0.32), Inches(box_w), Inches(0.5),
                 desc, size=9, color=DARK_GRAY, align=PP_ALIGN.CENTER)
        if i < len(boxes_top) - 1:
            arrow_x = x + box_w + 0.02
            add_arrow(slide, Inches(arrow_x), Inches(y + 0.32), Inches(0.21), Inches(0.2), fill=TEAL)

    # フェーズ2: クエリ
    add_text(slide, Inches(0.6), Inches(2.85), Inches(9.0), Inches(0.3),
             "フェーズ2: 質問応答", size=12, bold=True, color=TEAL)

    boxes_bot = [
        ("質問", "ユーザーの\n質問"),
        ("Embedding", "クエリの\nベクトル化"),
        ("検索", "ChromaDBから\nTop-8チャンク"),
        ("フィルタ", "距離0.45以下\nのみ採用"),
        ("LLM", "Gemini 2.5 Flash\nで回答生成"),
    ]
    y = 3.2
    for i, (label, desc) in enumerate(boxes_bot):
        x = start_x + i * (box_w + gap)
        add_round_rect(slide, Inches(x), Inches(y), Inches(box_w), Inches(box_h),
                       fill=PALE_CYAN, line=PURPLE_BLUE, line_width=0.75)
        add_text(slide, Inches(x), Inches(y + 0.05), Inches(box_w), Inches(0.3),
                 label, size=10, bold=True, color=PURPLE_BLUE, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(x), Inches(y + 0.32), Inches(box_w), Inches(0.5),
                 desc, size=9, color=DARK_GRAY, align=PP_ALIGN.CENTER)
        if i < len(boxes_bot) - 1:
            arrow_x = x + box_w + 0.02
            add_arrow(slide, Inches(arrow_x), Inches(y + 0.32), Inches(0.21), Inches(0.2), fill=PURPLE_BLUE)

    # 技術スタック
    add_round_rect(slide, Inches(0.6), Inches(4.4), Inches(8.8), Inches(0.7),
                   fill=LIGHT_GRAY)
    add_text(slide, Inches(0.85), Inches(4.5), Inches(8.5), Inches(0.3),
             "技術スタック", size=11, bold=True, color=NAVY)
    add_text(slide, Inches(0.85), Inches(4.75), Inches(8.5), Inches(0.3),
             "Python 3.11+ / ChromaDB / sentence-transformers (multilingual-e5-large) / Gemini API",
             size=10, color=DARK_GRAY)


# =====================================================
# Slide 4: ソース文書
# =====================================================
def slide_sources(slide, page, total):
    add_title_bar(slide, "3. ソース文書（10ファイル）")

    add_text(slide, Inches(0.6), Inches(1.1), Inches(9.0), Inches(0.4),
             "Anthropic 公式ドキュメント（platform.claude.com/docs/ja）から10ページを選定",
             size=12, color=DARK_GRAY)

    sources = [
        ("models_overview.md", "Claudeモデル一覧・料金・スペック比較表"),
        ("pricing.md", "詳細料金表（11モデル × 5列）"),
        ("tool_use.md", "Tool Use機能・関数呼び出し"),
        ("vision.md", "画像認識・対応フォーマット"),
        ("streaming.md", "ストリーミング・SSEイベント"),
        ("structured_outputs.md", "JSON出力・厳格なツール使用"),
        ("extended_thinking.md", "拡張思考・適応思考"),
        ("prompt_engineering.md", "プロンプトのベストプラクティス"),
        ("batch_processing.md", "バッチAPIの仕様"),
        ("pdf_support.md", "PDF処理機能"),
    ]

    col_w = 4.5
    row_h = 0.32
    for i, (name, desc) in enumerate(sources):
        col = i // 5
        row = i % 5
        x = 0.6 + col * (col_w + 0.2)
        y = 1.6 + row * (row_h + 0.05)
        add_rect(slide, Inches(x), Inches(y), Inches(0.05), Inches(row_h), fill=TEAL)
        add_text(slide, Inches(x + 0.12), Inches(y), Inches(1.8), Inches(row_h),
                 name, size=10, bold=True, color=NAVY)
        add_text(slide, Inches(x + 1.95), Inches(y), Inches(col_w - 1.95), Inches(row_h),
                 desc, size=9, color=DARK_GRAY)

    # ポイント
    add_round_rect(slide, Inches(0.6), Inches(3.6), Inches(8.8), Inches(1.4),
                   fill=PALE_CYAN)
    add_text(slide, Inches(0.85), Inches(3.7), Inches(8.4), Inches(0.3),
             "選定のポイント", size=12, bold=True, color=NAVY)
    points = [
        "・ 全て日本語の公式ドキュメント（クロスリンガル問題を回避）",
        "・ 数値・表・固有名詞が豊富で「正解が一意に決まる」",
        "・ AIが学習データから推測しないよう、最新モデル（Opus 4.6 等）の情報を含む",
        "・ NotebookLMにも同じファイルをアップロードして公平に比較できる",
    ]
    py = 4.0
    for p in points:
        add_text(slide, Inches(0.95), Inches(py), Inches(8.4), Inches(0.25),
                 p, size=10, color=BLACK)
        py += 0.23


# =====================================================
# Slide 5: 評価ルーブリックの進化
# =====================================================
def slide_rubric(slide, page, total):
    add_title_bar(slide, "4. 評価ルーブリックの設計と改善")

    add_text(slide, Inches(0.6), Inches(1.05), Inches(9.0), Inches(0.4),
             "教授に詰められないよう、採点基準を明確化（途中で気づいた問題点を反映）", size=12, color=DARK_GRAY)

    # v1
    add_round_rect(slide, Inches(0.6), Inches(1.55), Inches(2.85), Inches(3.4),
                   fill=PALE_CYAN, line=GRAY, line_width=0.5)
    add_text(slide, Inches(0.7), Inches(1.65), Inches(2.7), Inches(0.35),
             "v1: 主観評価", size=12, bold=True, color=NAVY)
    add_text(slide, Inches(0.7), Inches(1.95), Inches(2.7), Inches(0.3),
             "AI生成の期待回答との一致度", size=9, color=DARK_GRAY)
    add_text(slide, Inches(0.7), Inches(2.3), Inches(2.7), Inches(0.3),
             "問題点", size=10, bold=True, color=RED)
    add_text(slide, Inches(0.7), Inches(2.55), Inches(2.7), Inches(2.4),
             "・ 再現性なし\n・ 採点者バイアス\n・ 正解の妥当性が曖昧\n・ 「軽微な誤り」と「重大な誤り」の境界不明\n・ ポイントの粒度不明",
             size=9, color=BLACK)

    # 矢印
    add_arrow(slide, Inches(3.55), Inches(3.0), Inches(0.3), Inches(0.3), fill=TEAL)

    # v2
    add_round_rect(slide, Inches(3.95), Inches(1.55), Inches(2.85), Inches(3.4),
                   fill=LIGHT_CYAN, line=GRAY, line_width=0.5)
    add_text(slide, Inches(4.05), Inches(1.65), Inches(2.7), Inches(0.35),
             "v2: ルール詳細化", size=12, bold=True, color=NAVY)
    add_text(slide, Inches(4.05), Inches(1.95), Inches(2.7), Inches(0.3),
             "三段論法テスト・カバー率計算", size=9, color=DARK_GRAY)
    add_text(slide, Inches(4.05), Inches(2.3), Inches(2.7), Inches(0.3),
             "改善点", size=10, bold=True, color=GREEN)
    add_text(slide, Inches(4.05), Inches(2.55), Inches(2.7), Inches(2.4),
             "・ 平均スコア採用\n・ 推論明示=5/暗黙=3\n・ 古い情報=2\n・ 単語レベルでポイント抽出\n・ 迷ったら厳しい方",
             size=9, color=BLACK)

    # 矢印
    add_arrow(slide, Inches(6.9), Inches(3.0), Inches(0.3), Inches(0.3), fill=TEAL)

    # v3
    add_round_rect(slide, Inches(7.3), Inches(1.55), Inches(2.1), Inches(3.4),
                   fill=LIGHT_CYAN, line=NAVY, line_width=1.2)
    add_text(slide, Inches(7.4), Inches(1.65), Inches(1.95), Inches(0.35),
             "v3: 客観評価へ", size=12, bold=True, color=NAVY)
    add_text(slide, Inches(7.4), Inches(1.95), Inches(1.95), Inches(0.3),
             "正解一致型に転換", size=9, color=DARK_GRAY)
    add_text(slide, Inches(7.4), Inches(2.3), Inches(1.95), Inches(0.3),
             "決定版", size=10, bold=True, color=TEAL)
    add_text(slide, Inches(7.4), Inches(2.55), Inches(1.95), Inches(2.4),
             "・ 1問1答\n・ 単語/数値で\n　明確に正誤\n・ 模範回答と\n　ソース位置\n　を明示",
             size=9, color=BLACK)

    # 教訓
    add_text(slide, Inches(0.6), Inches(5.05), Inches(9.0), Inches(0.3),
             "学んだこと: 「正解が一意に決まる質問」を作ることが、再現性のある評価の前提条件",
             size=10, bold=True, color=TEAL)


# =====================================================
# Slide 6: 質問セットの進化
# =====================================================
def slide_question_evolution(slide, page, total):
    add_title_bar(slide, "5. 質問セットの進化（v1〜v6）")

    rows = [
        ("v1", "汎用質問", "主観", "「Constitutional AIとは何ですか？」", "曖昧", LIGHT_CYAN),
        ("v2", "翻訳実験", "主観", "（v1と同じ質問・クエリ翻訳実験）", "悪化", PALE_CYAN),
        ("v3", "フィルタ実験", "主観", "（v1と同じ質問・距離フィルタ追加）", "改善", PALE_CYAN),
        ("v4", "事実抽出", "客観", "「Haiku 4.5の入力料金は？」 → $1", "高得点", LIGHT_CYAN),
        ("v5", "マルチソース", "客観", "「100ページPDFをBatchで処理した料金は？」", "天井", LIGHT_CYAN),
        ("v6", "RAG弱点特化", "客観", "「適応思考: いいえのモデルは何個？」", "失敗発生", LIME),
    ]

    # ヘッダー
    add_rect(slide, Inches(0.6), Inches(1.15), Inches(0.7), Inches(0.4), fill=NAVY)
    add_rect(slide, Inches(1.3), Inches(1.15), Inches(1.5), Inches(0.4), fill=NAVY)
    add_rect(slide, Inches(2.8), Inches(1.15), Inches(1.0), Inches(0.4), fill=NAVY)
    add_rect(slide, Inches(3.8), Inches(1.15), Inches(4.5), Inches(0.4), fill=NAVY)
    add_rect(slide, Inches(8.3), Inches(1.15), Inches(1.1), Inches(0.4), fill=NAVY)
    headers = ["バージョン", "テーマ", "評価方式", "質問例", "結果"]
    cols_x = [0.6, 1.3, 2.8, 3.8, 8.3]
    cols_w = [0.7, 1.5, 1.0, 4.5, 1.1]
    for i, h in enumerate(headers):
        add_text(slide, Inches(cols_x[i]), Inches(1.18), Inches(cols_w[i]), Inches(0.35),
                 h, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 行
    y = 1.55
    row_h = 0.5
    for ver, theme, eval_type, example, result, color in rows:
        add_rect(slide, Inches(0.6), Inches(y), Inches(0.7), Inches(row_h), fill=color)
        add_rect(slide, Inches(1.3), Inches(y), Inches(1.5), Inches(row_h), fill=color)
        add_rect(slide, Inches(2.8), Inches(y), Inches(1.0), Inches(row_h), fill=color)
        add_rect(slide, Inches(3.8), Inches(y), Inches(4.5), Inches(row_h), fill=color)
        add_rect(slide, Inches(8.3), Inches(y), Inches(1.1), Inches(row_h), fill=color)

        add_text(slide, Inches(0.6), Inches(y + 0.13), Inches(0.7), Inches(0.3),
                 ver, size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(1.4), Inches(y + 0.13), Inches(1.4), Inches(0.3),
                 theme, size=10, color=BLACK)
        add_text(slide, Inches(2.8), Inches(y + 0.13), Inches(1.0), Inches(0.3),
                 eval_type, size=10, color=BLACK, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(3.9), Inches(y + 0.13), Inches(4.4), Inches(0.3),
                 example, size=9, color=DARK_GRAY)
        add_text(slide, Inches(8.3), Inches(y + 0.13), Inches(1.1), Inches(0.3),
                 result, size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        y += row_h

    add_text(slide, Inches(0.6), Inches(4.75), Inches(9.0), Inches(0.3),
             "気づき", size=11, bold=True, color=TEAL)
    add_text(slide, Inches(0.6), Inches(5.0), Inches(9.0), Inches(0.4),
             "v5で天井（100%）に達した後、RAGの弱点を意図的に突くv6を作成。これが転機となった。",
             size=10, color=BLACK)


# =====================================================
# Slide 7: スコア推移グラフ
# =====================================================
def slide_score_chart(slide, page, total):
    add_title_bar(slide, "6. スコア推移（v1 → v7）")

    # グラフエリア
    chart_left = 1.0
    chart_top = 1.4
    chart_w = 8.0
    chart_h = 3.0

    # 枠
    add_rect(slide, Inches(chart_left), Inches(chart_top), Inches(chart_w), Inches(chart_h),
             fill=WHITE, line=GRAY, line_width=0.75)

    # Y軸ラベル
    for v, lbl in [(0, "0%"), (0.25, "25%"), (0.5, "50%"), (0.75, "75%"), (1.0, "100%")]:
        y = chart_top + chart_h * (1 - v)
        add_text(slide, Inches(chart_left - 0.55), Inches(y - 0.1), Inches(0.5), Inches(0.2),
                 lbl, size=8, color=GRAY, align=PP_ALIGN.RIGHT)
        add_line(slide, Inches(chart_left), Inches(y), Inches(chart_left + chart_w), Inches(y),
                 color=LIGHT_GRAY, weight=0.5)

    # データ（v1〜v7のRAGスコア + NotebookLM比較）
    scores = [
        ("v1", 0.84, "baseline"),  # 4.2/5
        ("v2", 0.78, "翻訳"),  # 3.9/5
        ("v3a", 0.88, "フィルタ"),  # 4.4/5
        ("v4", 0.94, "事実抽出"),  # 47/50
        ("v5", 1.00, "マルチ"),  # 50/50
        ("v6", 0.80, "弱点特化"),  # 8/10
        ("v7", 0.90, "表認識"),  # 9/10
    ]
    n = len(scores)
    bar_area_w = chart_w - 1.0
    bar_w = bar_area_w / n * 0.55
    bar_gap = bar_area_w / n

    for i, (ver, score, label) in enumerate(scores):
        x = chart_left + 0.5 + i * bar_gap
        bar_h_actual = chart_h * score
        bar_y = chart_top + chart_h - bar_h_actual

        # 色（v6は赤系、v7は緑系で強調）
        if ver == "v6":
            bar_color = ORANGE
        elif ver == "v7":
            bar_color = LIME
        else:
            bar_color = TEAL

        add_rect(slide, Inches(x), Inches(bar_y), Inches(bar_w), Inches(bar_h_actual),
                 fill=bar_color)

        # スコア値
        add_text(slide, Inches(x - 0.1), Inches(bar_y - 0.25), Inches(bar_w + 0.2), Inches(0.2),
                 f"{int(score*100)}%", size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

        # ラベル
        add_text(slide, Inches(x - 0.2), Inches(chart_top + chart_h + 0.05), Inches(bar_w + 0.4), Inches(0.2),
                 ver, size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(x - 0.3), Inches(chart_top + chart_h + 0.25), Inches(bar_w + 0.6), Inches(0.2),
                 label, size=8, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # NotebookLM参照線
    nlm_y = chart_top + chart_h * (1 - 1.0)
    add_line(slide, Inches(chart_left), Inches(nlm_y), Inches(chart_left + chart_w), Inches(nlm_y),
             color=PURPLE_BLUE, weight=2.0)
    add_text(slide, Inches(chart_left + chart_w - 1.8), Inches(nlm_y - 0.18), Inches(1.7), Inches(0.2),
             "NotebookLM (100%)", size=8, bold=True, color=PURPLE_BLUE, align=PP_ALIGN.RIGHT)

    # コメント
    add_text(slide, Inches(0.6), Inches(4.7), Inches(9.0), Inches(0.3),
             "観察", size=11, bold=True, color=TEAL)
    add_text(slide, Inches(0.6), Inches(4.95), Inches(9.0), Inches(0.4),
             "v6でRAGの弱点が露呈（80%）→ v7で表認識チャンキングを実装し90%まで改善",
             size=10, color=BLACK)


# =====================================================
# Slide 8: v6 質問例
# =====================================================
def slide_v6_questions(slide, page, total):
    add_title_bar(slide, "7. v6: RAGの弱点を突いた10問")

    add_text(slide, Inches(0.6), Inches(1.05), Inches(9.0), Inches(0.4),
             "RAGの構造的弱点（チャンク分割によって失われる情報）を意図的にテスト",
             size=11, color=DARK_GRAY)

    questions = [
        ("Q1", "否定", "Haikuだけ1Mトークン非対応", "Claude Haiku 4.5"),
        ("Q2", "カウント", "「適応思考: いいえ」のモデル数", "2個"),
        ("Q3", "紛らわしい命名", "Claude 3 Haikuの入力料金", "$0.25"),
        ("Q4", "順位付け", "最安の出力料金モデル", "Claude Haiku 3"),
        ("Q5", "記載なし推論", "リストにOpus 3はあるか", "いいえ"),
        ("Q6", "複数セル算術", "1hキャッシュ料金の差額", "$4"),
        ("Q7", "バージョン混同", "JSON出力の最新パラメータ名", "output_config.format"),
        ("Q8", "比率・乗数", "5分キャッシュは標準の何倍", "1.25倍"),
        ("Q9", "暗黙の制約", "200ページPDFは処理可能？", "いいえ（100上限）"),
        ("Q10", "非推奨判定", "Opus 4.6で非推奨のtype値", "\"enabled\""),
    ]

    y = 1.55
    row_h = 0.32
    # ヘッダー
    add_rect(slide, Inches(0.6), Inches(y), Inches(0.5), Inches(row_h), fill=NAVY)
    add_rect(slide, Inches(1.1), Inches(y), Inches(2.0), Inches(row_h), fill=NAVY)
    add_rect(slide, Inches(3.1), Inches(y), Inches(4.5), Inches(row_h), fill=NAVY)
    add_rect(slide, Inches(7.6), Inches(y), Inches(1.8), Inches(row_h), fill=NAVY)
    headers = ["#", "RAG弱点カテゴリ", "質問内容", "正解"]
    xs = [0.6, 1.1, 3.1, 7.6]
    ws = [0.5, 2.0, 4.5, 1.8]
    for i, h in enumerate(headers):
        add_text(slide, Inches(xs[i]), Inches(y + 0.05), Inches(ws[i]), Inches(0.25),
                 h, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    y += row_h

    for q_id, cat, content, ans in questions:
        is_failed = q_id in ["Q2", "Q6"]
        bg = LIME if is_failed else PALE_CYAN
        add_rect(slide, Inches(0.6), Inches(y), Inches(0.5), Inches(row_h), fill=bg)
        add_rect(slide, Inches(1.1), Inches(y), Inches(2.0), Inches(row_h), fill=bg)
        add_rect(slide, Inches(3.1), Inches(y), Inches(4.5), Inches(row_h), fill=bg)
        add_rect(slide, Inches(7.6), Inches(y), Inches(1.8), Inches(row_h), fill=bg)

        add_text(slide, Inches(0.6), Inches(y + 0.06), Inches(0.5), Inches(0.25),
                 q_id, size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(1.2), Inches(y + 0.06), Inches(1.9), Inches(0.25),
                 cat, size=8, color=BLACK)
        add_text(slide, Inches(3.2), Inches(y + 0.06), Inches(4.4), Inches(0.25),
                 content, size=8, color=BLACK)
        add_text(slide, Inches(7.6), Inches(y + 0.06), Inches(1.8), Inches(0.25),
                 ans, size=8, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        y += row_h

    add_text(slide, Inches(0.6), Inches(5.0), Inches(9.0), Inches(0.3),
             "黄緑のQ2・Q6でRAGが失敗（NotebookLMは正解）",
             size=10, bold=True, color=ORANGE)


# =====================================================
# Slide 9: v6 結果ハイライト
# =====================================================
def slide_v6_results(slide, page, total):
    add_title_bar(slide, "8. v6 結果: RAG vs NotebookLM 直接比較")

    # 大きな比較
    rag_x = 0.8
    nlm_x = 5.4
    box_w = 4.0
    box_h = 2.2

    # RAG
    add_round_rect(slide, Inches(rag_x), Inches(1.2), Inches(box_w), Inches(box_h),
                   fill=LIGHT_CYAN, line=NAVY, line_width=1.2)
    add_text(slide, Inches(rag_x + 0.2), Inches(1.3), Inches(box_w - 0.4), Inches(0.4),
             "自作 RAG (v6)", size=14, bold=True, color=NAVY)
    add_text(slide, Inches(rag_x + 0.2), Inches(1.7), Inches(box_w - 0.4), Inches(1.0),
             "8 / 10", size=48, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(rag_x + 0.2), Inches(2.7), Inches(box_w - 0.4), Inches(0.3),
             "正答率 80%", size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(rag_x + 0.2), Inches(3.05), Inches(box_w - 0.4), Inches(0.3),
             "失敗: Q2（カウント）, Q6（複数セル算術）",
             size=10, color=ORANGE, align=PP_ALIGN.CENTER)

    # NotebookLM
    add_round_rect(slide, Inches(nlm_x), Inches(1.2), Inches(box_w), Inches(box_h),
                   fill=PALE_CYAN, line=PURPLE_BLUE, line_width=1.2)
    add_text(slide, Inches(nlm_x + 0.2), Inches(1.3), Inches(box_w - 0.4), Inches(0.4),
             "Google NotebookLM", size=14, bold=True, color=PURPLE_BLUE)
    add_text(slide, Inches(nlm_x + 0.2), Inches(1.7), Inches(box_w - 0.4), Inches(1.0),
             "10 / 10", size=48, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(nlm_x + 0.2), Inches(2.7), Inches(box_w - 0.4), Inches(0.3),
             "正答率 100%", size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(nlm_x + 0.2), Inches(3.05), Inches(box_w - 0.4), Inches(0.3),
             "全問正解", size=10, color=GREEN, align=PP_ALIGN.CENTER)

    # 決定的な差
    add_round_rect(slide, Inches(0.6), Inches(3.7), Inches(8.8), Inches(1.5),
                   fill=LIGHT_GRAY)
    add_text(slide, Inches(0.85), Inches(3.8), Inches(8.4), Inches(0.3),
             "決定的な差: 表データに対する処理能力", size=12, bold=True, color=NAVY)
    add_text(slide, Inches(0.85), Inches(4.1), Inches(8.4), Inches(0.3),
             "Q2 (カウント): RAG「情報がない」 ↔ NLM「2つ（Sonnet 4.5とHaiku 4.5）」",
             size=10, color=BLACK)
    add_text(slide, Inches(0.85), Inches(4.35), Inches(8.4), Inches(0.3),
             "Q6 (算術): RAG「期間の区別なし」（誤答） ↔ NLM「$4 = $6 - $2」",
             size=10, color=BLACK)
    add_text(slide, Inches(0.85), Inches(4.7), Inches(8.4), Inches(0.4),
             "→ NotebookLMは表全体を保持しているため、行・列をまたぐ操作ができる",
             size=11, bold=True, color=TEAL)


# =====================================================
# Slide 10: なぜNotebookLMが正解できるか
# =====================================================
def slide_why_nlm_works(slide, page, total):
    add_title_bar(slide, "9. なぜNotebookLMは正解できるのか？")

    add_text(slide, Inches(0.6), Inches(1.05), Inches(9.0), Inches(0.3),
             "RAG と NotebookLM はアーキテクチャが構造的に異なる",
             size=12, color=DARK_GRAY)

    # RAGの図
    add_text(slide, Inches(0.6), Inches(1.45), Inches(4.4), Inches(0.3),
             "自作RAGの場合", size=12, bold=True, color=ORANGE)

    # 表の分割イメージ
    add_rect(slide, Inches(0.6), Inches(1.8), Inches(4.4), Inches(0.4),
             fill=NAVY)
    add_text(slide, Inches(0.6), Inches(1.85), Inches(4.4), Inches(0.3),
             "| 機能 | Opus 4.6 | Sonnet 4.5 | Haiku 4.5 |",
             size=8, color=WHITE, align=PP_ALIGN.CENTER)

    # 切断線
    add_text(slide, Inches(0.6), Inches(2.25), Inches(4.4), Inches(0.2),
             "─ ─ ─ ─ ─ チャンク境界 ─ ─ ─ ─ ─",
             size=7, color=RED, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0.6), Inches(2.5), Inches(4.4), Inches(0.4),
             fill=PALE_CYAN, line=GRAY, line_width=0.5)
    add_text(slide, Inches(0.6), Inches(2.55), Inches(4.4), Inches(0.3),
             "| 適応思考 | はい | いいえ | いいえ |",
             size=8, color=BLACK, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.6), Inches(3.0), Inches(4.4), Inches(0.4),
             "→ ヘッダーと切り離されると\n　 「いいえ」がどのモデルか不明",
             size=10, color=ORANGE)

    # NotebookLMの図
    add_text(slide, Inches(5.2), Inches(1.45), Inches(4.4), Inches(0.3),
             "NotebookLM の場合", size=12, bold=True, color=GREEN)

    add_round_rect(slide, Inches(5.2), Inches(1.8), Inches(4.4), Inches(1.5),
                   fill=PALE_CYAN, line=PURPLE_BLUE, line_width=1.0)
    add_text(slide, Inches(5.3), Inches(1.85), Inches(4.2), Inches(0.25),
             "| 機能 | Opus 4.6 | Sonnet 4.5 | Haiku 4.5 |",
             size=8, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(5.3), Inches(2.1), Inches(4.2), Inches(0.25),
             "| 料金 | $5 | $3 | $1 |",
             size=8, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(5.3), Inches(2.35), Inches(4.2), Inches(0.25),
             "| 拡張思考 | はい | はい | はい |",
             size=8, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(5.3), Inches(2.6), Inches(4.2), Inches(0.25),
             "| 適応思考 | はい | いいえ | いいえ |",
             size=8, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(5.3), Inches(2.9), Inches(4.2), Inches(0.3),
             "（表全体を1つの単位として保持）",
             size=8, color=GRAY, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(5.2), Inches(3.0), Inches(4.4), Inches(0.4),
             "→ 行・列の対応が完全に保たれ\n　 全体構造を理解した回答が可能",
             size=10, color=GREEN)

    # 結論
    add_round_rect(slide, Inches(0.6), Inches(4.5), Inches(8.8), Inches(0.8),
                   fill=LIGHT_CYAN)
    add_text(slide, Inches(0.85), Inches(4.6), Inches(8.4), Inches(0.3),
             "本質的な違い", size=11, bold=True, color=NAVY)
    add_text(slide, Inches(0.85), Inches(4.85), Inches(8.4), Inches(0.4),
             "RAG = 文書を分割して必要な部分だけ見る　／　NotebookLM = 文書全体を保持して全体構造を把握",
             size=10, color=BLACK)


# =====================================================
# Slide 11: v7 改善実装
# =====================================================
def slide_v7_implementation(slide, page, total):
    add_title_bar(slide, "10. v7: 表認識チャンキングの実装")

    # 改善のアイデア
    add_text(slide, Inches(0.6), Inches(1.05), Inches(9.0), Inches(0.4),
             "改善案: Markdownの表を検出し、表全体を1つのチャンクとして保持する",
             size=12, color=DARK_GRAY)

    # Before
    add_text(slide, Inches(0.6), Inches(1.55), Inches(4.4), Inches(0.3),
             "Before (v6): 単純チャンク分割", size=11, bold=True, color=ORANGE)

    add_rect(slide, Inches(0.6), Inches(1.9), Inches(4.4), Inches(0.4),
             fill=NAVY)
    add_text(slide, Inches(0.6), Inches(1.95), Inches(4.4), Inches(0.3),
             "ヘッダー行（500文字制限で別チャンク）",
             size=9, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.6), Inches(2.35), Inches(4.4), Inches(0.2),
             "✕ 切断", size=8, color=RED, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0.6), Inches(2.55), Inches(4.4), Inches(0.4),
             fill=PALE_CYAN)
    add_text(slide, Inches(0.6), Inches(2.6), Inches(4.4), Inches(0.3),
             "データ行A・B（独立チャンク）",
             size=9, color=BLACK, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0.6), Inches(3.0), Inches(4.4), Inches(0.4),
             fill=PALE_CYAN)
    add_text(slide, Inches(0.6), Inches(3.05), Inches(4.4), Inches(0.3),
             "データ行C・D（独立チャンク）",
             size=9, color=BLACK, align=PP_ALIGN.CENTER)

    # After
    add_text(slide, Inches(5.2), Inches(1.55), Inches(4.4), Inches(0.3),
             "After (v7): 表認識チャンキング", size=11, bold=True, color=GREEN)

    add_round_rect(slide, Inches(5.2), Inches(1.9), Inches(4.4), Inches(1.5),
                   fill=PALE_CYAN, line=GREEN, line_width=1.5)
    add_text(slide, Inches(5.3), Inches(1.95), Inches(4.2), Inches(0.25),
             "ヘッダー行", size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(5.3), Inches(2.2), Inches(4.2), Inches(0.25),
             "データ行A", size=9, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(5.3), Inches(2.4), Inches(4.2), Inches(0.25),
             "データ行B", size=9, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(5.3), Inches(2.6), Inches(4.2), Inches(0.25),
             "データ行C", size=9, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(5.3), Inches(2.8), Inches(4.2), Inches(0.25),
             "データ行D", size=9, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(5.3), Inches(3.05), Inches(4.2), Inches(0.3),
             "（全行が1つのチャンクに）",
             size=8, color=GRAY, align=PP_ALIGN.CENTER)

    # 実装の核心
    add_round_rect(slide, Inches(0.6), Inches(3.85), Inches(8.8), Inches(1.4),
                   fill=LIGHT_GRAY)
    add_text(slide, Inches(0.85), Inches(3.95), Inches(8.4), Inches(0.3),
             "実装の核心 (src/chunker.py)", size=11, bold=True, color=NAVY)
    add_text(slide, Inches(0.85), Inches(4.25), Inches(8.4), Inches(0.3),
             "1. テキストを「表ブロック」と「通常テキストブロック」に分解",
             size=10, color=BLACK)
    add_text(slide, Inches(0.85), Inches(4.5), Inches(8.4), Inches(0.3),
             "2. Markdownの`|`で始まる連続行を検出し、表として1チャンクにまとめる",
             size=10, color=BLACK)
    add_text(slide, Inches(0.85), Inches(4.75), Inches(8.4), Inches(0.3),
             "3. 表の直前のセクション見出し（##）も同じチャンクに含める",
             size=10, color=BLACK)
    add_text(slide, Inches(0.85), Inches(5.0), Inches(8.4), Inches(0.3),
             "4. 各チャンクに metadata.block_type で「table」「text」のラベルを付与",
             size=10, color=BLACK)


# =====================================================
# Slide 12: v7 結果
# =====================================================
def slide_v7_results(slide, page, total):
    add_title_bar(slide, "11. v7 結果: 9/10（+1問改善）")

    # スコア比較
    add_text(slide, Inches(0.6), Inches(1.1), Inches(9.0), Inches(0.3),
             "v6 → v7 の質問別比較",
             size=12, bold=True, color=TEAL)

    # ヘッダー
    y = 1.5
    add_rect(slide, Inches(0.6), Inches(y), Inches(0.6), Inches(0.35), fill=NAVY)
    add_rect(slide, Inches(1.2), Inches(y), Inches(4.0), Inches(0.35), fill=NAVY)
    add_rect(slide, Inches(5.2), Inches(y), Inches(2.1), Inches(0.35), fill=NAVY)
    add_rect(slide, Inches(7.3), Inches(y), Inches(1.05), Inches(0.35), fill=NAVY)
    add_rect(slide, Inches(8.35), Inches(y), Inches(1.05), Inches(0.35), fill=NAVY)
    headers = ["#", "カテゴリ", "正解", "v6", "v7"]
    xs = [0.6, 1.2, 5.2, 7.3, 8.35]
    ws = [0.6, 4.0, 2.1, 1.05, 1.05]
    for i, h in enumerate(headers):
        add_text(slide, Inches(xs[i]), Inches(y + 0.05), Inches(ws[i]), Inches(0.25),
                 h, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    questions = [
        ("Q1", "否定（〜しない）", "Haiku 4.5", "○", "○", PALE_CYAN),
        ("Q2", "カウント", "2個", "✕", "✕", LIGHT_CYAN),
        ("Q3", "紛らわしい命名", "$0.25", "○", "○", PALE_CYAN),
        ("Q4", "順位付け", "Haiku 3", "○", "○", PALE_CYAN),
        ("Q5", "記載なし推論", "いいえ", "○", "○", PALE_CYAN),
        ("Q6", "複数セル算術", "$4", "✕", "○", LIME),
        ("Q7", "バージョン混同", "output_config.format", "○", "○", PALE_CYAN),
        ("Q8", "比率・乗数", "1.25倍", "○", "○", PALE_CYAN),
        ("Q9", "暗黙の制約判定", "いいえ", "○", "○", PALE_CYAN),
        ("Q10", "非推奨判定", "\"enabled\"", "○", "○", PALE_CYAN),
    ]
    y += 0.35
    row_h = 0.27
    for q_id, cat, ans, v6, v7, color in questions:
        add_rect(slide, Inches(0.6), Inches(y), Inches(0.6), Inches(row_h), fill=color)
        add_rect(slide, Inches(1.2), Inches(y), Inches(4.0), Inches(row_h), fill=color)
        add_rect(slide, Inches(5.2), Inches(y), Inches(2.1), Inches(row_h), fill=color)
        add_rect(slide, Inches(7.3), Inches(y), Inches(1.05), Inches(row_h), fill=color)
        add_rect(slide, Inches(8.35), Inches(y), Inches(1.05), Inches(row_h), fill=color)

        add_text(slide, Inches(0.6), Inches(y + 0.04), Inches(0.6), Inches(0.2),
                 q_id, size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(1.3), Inches(y + 0.04), Inches(3.9), Inches(0.2),
                 cat, size=9, color=BLACK)
        add_text(slide, Inches(5.2), Inches(y + 0.04), Inches(2.1), Inches(0.2),
                 ans, size=9, color=BLACK, align=PP_ALIGN.CENTER)
        v6_color = GREEN if v6 == "○" else RED
        v7_color = GREEN if v7 == "○" else RED
        add_text(slide, Inches(7.3), Inches(y + 0.04), Inches(1.05), Inches(0.2),
                 v6, size=11, bold=True, color=v6_color, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(8.35), Inches(y + 0.04), Inches(1.05), Inches(0.2),
                 v7, size=11, bold=True, color=v7_color, align=PP_ALIGN.CENTER)
        y += row_h

    # 注釈
    add_text(slide, Inches(0.6), Inches(5.05), Inches(9.0), Inches(0.3),
             "Q6が解けるようになったのは、料金表全体（11モデル × 5列）が1チャンクで取得されたため",
             size=9, color=DARK_GRAY)


# =====================================================
# Slide 13: Q2 が解けない理由
# =====================================================
def slide_q2_analysis(slide, page, total):
    add_title_bar(slide, "12. v7でもQ2が解けない理由")

    add_text(slide, Inches(0.6), Inches(1.05), Inches(9.0), Inches(0.3),
             "表認識チャンキングだけでは解決しない、別の問題が判明",
             size=12, color=DARK_GRAY)

    # 問題の説明
    add_round_rect(slide, Inches(0.6), Inches(1.5), Inches(8.8), Inches(0.7), fill=PALE_CYAN)
    add_text(slide, Inches(0.85), Inches(1.6), Inches(8.4), Inches(0.3),
             "質問: 「最新モデル比較表で『適応思考: いいえ』のモデルは何個？」",
             size=11, bold=True, color=NAVY)
    add_text(slide, Inches(0.85), Inches(1.85), Inches(8.4), Inches(0.3),
             "→ RAGの回答: 「資料に情報が含まれていません」（v6・v7とも同じ）",
             size=10, color=ORANGE)

    # 検索結果のランキング
    add_text(slide, Inches(0.6), Inches(2.4), Inches(9.0), Inches(0.3),
             "v7の検索結果 Top 8 を確認すると…",
             size=11, bold=True, color=TEAL)

    rankings = [
        ("1位", "extended-thinking.md (説明文)", 0.328, False),
        ("2位", "prompt-engineering.md (説明文)", 0.329, False),
        ("3位", "extended-thinking.md (説明文)", 0.352, False),
        ("4位", "prompt-engineering.md (説明文)", 0.353, False),
        ("5〜7位", "extended-thinking.md (説明文)", 0.356, False),
        ("8位", "extended-thinking.md (表)", 0.366, False),
        ("圏外", "models_overview.md (表) ← 必要なチャンク", 0.40, True),
    ]
    y = 2.75
    for rank, source, dist, is_target in rankings:
        bg = LIME if is_target else PALE_CYAN
        add_rect(slide, Inches(0.6), Inches(y), Inches(8.8), Inches(0.27), fill=bg)
        add_text(slide, Inches(0.7), Inches(y + 0.04), Inches(0.9), Inches(0.2),
                 rank, size=9, bold=True, color=NAVY)
        add_text(slide, Inches(1.6), Inches(y + 0.04), Inches(6.0), Inches(0.2),
                 source, size=9, color=BLACK)
        add_text(slide, Inches(7.6), Inches(y + 0.04), Inches(1.7), Inches(0.2),
                 f"距離 {dist:.3f}", size=9, color=GRAY, align=PP_ALIGN.RIGHT)
        y += 0.28

    add_round_rect(slide, Inches(0.6), Inches(4.85), Inches(8.8), Inches(0.55), fill=LIGHT_GRAY)
    add_text(slide, Inches(0.85), Inches(4.92), Inches(8.4), Inches(0.3),
             "原因: 「適応思考」キーワードが説明文に多数出現するため、説明文チャンクが上位を占める",
             size=10, bold=True, color=NAVY)
    add_text(slide, Inches(0.85), Inches(5.15), Inches(8.4), Inches(0.3),
             "models_overview.mdの表内に「適応思考」は1行しかなく、相対的な類似度が低くなる",
             size=10, color=BLACK)


# =====================================================
# Slide 14: 結論
# =====================================================
def slide_conclusion(slide, page, total):
    add_title_bar(slide, "13. 結論")

    points = [
        ("①", "単純な事実抽出ではRAGとNotebookLMはほぼ同等",
         "v4（事実抽出）では RAG 94%、v5（マルチソース）では RAG 100%"),
        ("②", "RAGの構造的弱点は表データで露呈",
         "v6では RAG 80% vs NotebookLM 100%、Q2（カウント）とQ6（複数セル算術）で失敗"),
        ("③", "本質的な違いは「文書の渡し方」",
         "RAG = 部分（チャンク）／ NotebookLM = 全体保持。LLMの能力差ではない"),
        ("④", "工夫すればRAGも改善できる（v7で実証）",
         "表認識チャンキングで Q6 が解けるようになり、80% → 90% に向上"),
        ("⑤", "ただし汎用的な解決には限界がある",
         "Q2のような検索バイアスは別の改善が必要で、副作用なしの解決は難しい"),
    ]

    y = 1.2
    for num, headline, detail in points:
        add_round_rect(slide, Inches(0.6), Inches(y), Inches(8.8), Inches(0.78),
                       fill=PALE_CYAN, line=TEAL, line_width=0.5)
        add_text(slide, Inches(0.75), Inches(y + 0.12), Inches(0.5), Inches(0.4),
                 num, size=18, bold=True, color=TEAL)
        add_text(slide, Inches(1.3), Inches(y + 0.1), Inches(8.0), Inches(0.3),
                 headline, size=12, bold=True, color=NAVY)
        add_text(slide, Inches(1.3), Inches(y + 0.4), Inches(8.0), Inches(0.3),
                 detail, size=10, color=DARK_GRAY)
        y += 0.85


# =====================================================
# Slide 15: 学んだこと
# =====================================================
def slide_lessons(slide, page, total):
    add_title_bar(slide, "14. 実験を通じて学んだこと")

    # 左: 技術的な学び
    add_text(slide, Inches(0.6), Inches(1.1), Inches(4.4), Inches(0.4),
             "技術的な学び", size=14, bold=True, color=TEAL)

    tech_items = [
        ("チャンク分割は表データに不利",
         "ヘッダーとセルが分離すると意味が失われる"),
        ("クロスリンガル検索の限界",
         "all-MiniLM-L6-v2 → multilingual-e5-large に変更で大幅改善"),
        ("検索バイアスの存在",
         "頻出キーワードを含む説明文が上位を独占する"),
        ("プロンプトで挙動を制御可能",
         "推測禁止・論理的推論の明示で精度向上"),
    ]
    y = 1.55
    for title, desc in tech_items:
        add_rect(slide, Inches(0.6), Inches(y), Inches(0.08), Inches(0.6), fill=TEAL)
        add_text(slide, Inches(0.78), Inches(y), Inches(4.2), Inches(0.3),
                 title, size=10, bold=True, color=NAVY)
        add_text(slide, Inches(0.78), Inches(y + 0.25), Inches(4.2), Inches(0.3),
                 desc, size=9, color=DARK_GRAY)
        y += 0.7

    # 右: 実験設計の学び
    add_text(slide, Inches(5.2), Inches(1.1), Inches(4.4), Inches(0.4),
             "実験設計の学び", size=14, bold=True, color=TEAL)

    design_items = [
        ("採点基準の重要性",
         "再現性のある評価には客観的なルーブリックが必須"),
        ("質問設計が成果を左右する",
         "v5で天井→v6で意図的に弱点を突き差を発見"),
        ("ベースライン比較の価値",
         "NotebookLM比較で構造的な違いが明確化"),
        ("仮説→実装→検証のサイクル",
         "v7で仮説を実装し、改善と限界を同時に発見"),
    ]
    y = 1.55
    for title, desc in design_items:
        add_rect(slide, Inches(5.2), Inches(y), Inches(0.08), Inches(0.6), fill=PURPLE_BLUE)
        add_text(slide, Inches(5.38), Inches(y), Inches(4.2), Inches(0.3),
                 title, size=10, bold=True, color=NAVY)
        add_text(slide, Inches(5.38), Inches(y + 0.25), Inches(4.2), Inches(0.3),
                 desc, size=9, color=DARK_GRAY)
        y += 0.7

    # 最も重要な気づき
    add_round_rect(slide, Inches(0.6), Inches(4.5), Inches(8.8), Inches(0.7),
                   fill=LIGHT_CYAN, line=NAVY, line_width=1.0)
    add_text(slide, Inches(0.85), Inches(4.6), Inches(8.4), Inches(0.3),
             "最も重要な気づき", size=11, bold=True, color=NAVY)
    add_text(slide, Inches(0.85), Inches(4.85), Inches(8.4), Inches(0.3),
             "RAGの限界はLLMの能力ではなく、「情報の渡し方」に起因する構造的なもの",
             size=11, bold=True, color=TEAL)


# =====================================================
# Slide 16: 今後の課題
# =====================================================
def slide_future_work(slide, page, total):
    add_title_bar(slide, "15. 今後の課題")

    items = [
        ("Q2の解決",
         "「カウント」「集計」のような全行走査が必要な質問への対応",
         "ハイブリッド検索（セマンティック + BM25）/ メタデータフィルタリング"),
        ("検索バイアスの緩和",
         "頻出キーワードによる説明文偏重を解消",
         "Multi-Query Retrieval / リランカーの導入"),
        ("評価の自動化",
         "現在は人力での採点。再現性をさらに高めるには",
         "LLM-as-Judge による自動採点（評価専用エージェント）"),
        ("質問セットの拡充",
         "n=10では統計的有意性の議論ができない",
         "質問の複雑度別に各20問程度に拡張"),
        ("ハイブリッドアプローチ",
         "「RAG + 全文書投入」の組み合わせ",
         "検索結果に加えて文書の目次・見出しも提示"),
    ]

    y = 1.1
    for title, problem, solution in items:
        add_round_rect(slide, Inches(0.6), Inches(y), Inches(8.8), Inches(0.78),
                       fill=PALE_CYAN, line=GRAY, line_width=0.4)
        add_text(slide, Inches(0.85), Inches(y + 0.08), Inches(3.0), Inches(0.3),
                 title, size=12, bold=True, color=NAVY)
        add_text(slide, Inches(0.85), Inches(y + 0.35), Inches(8.4), Inches(0.25),
                 f"課題: {problem}", size=9, color=DARK_GRAY)
        add_text(slide, Inches(0.85), Inches(y + 0.55), Inches(8.4), Inches(0.25),
                 f"対策案: {solution}", size=9, color=TEAL)
        y += 0.85


# =====================================================
# Slide 17: 参考情報
# =====================================================
def slide_references(slide, page, total):
    add_title_bar(slide, "16. 参考情報・成果物")

    # GitHubリポジトリ
    add_text(slide, Inches(0.6), Inches(1.1), Inches(9.0), Inches(0.4),
             "ソースコード", size=14, bold=True, color=TEAL)
    add_round_rect(slide, Inches(0.6), Inches(1.5), Inches(8.8), Inches(0.55), fill=PALE_CYAN)
    add_text(slide, Inches(0.85), Inches(1.62), Inches(8.4), Inches(0.3),
             "GitHub: https://github.com/shogonakamura1/rag-comparison",
             size=11, color=NAVY)

    # 評価結果ファイル
    add_text(slide, Inches(0.6), Inches(2.2), Inches(9.0), Inches(0.4),
             "評価結果ファイル（evaluation/results/）", size=14, bold=True, color=TEAL)

    files = [
        ("v1_baseline_*.{json,md}", "ベースライン評価（旧質問セット）"),
        ("v2_query_translation_*", "クエリ翻訳実験"),
        ("v3_filtering_*", "距離フィルタリング + A/Bテスト"),
        ("v4_fact_extraction_*", "単一ソース事実抽出（94%）"),
        ("v5_multi_source_*", "複数ソース・計算問題（100%）"),
        ("v6_weakness_targeted_*", "RAG弱点特化（80%）"),
        ("v6_rag_vs_notebooklm_comparison.md", "NotebookLM 直接比較レポート"),
        ("v7_table_chunking_*", "表認識チャンキング（90%）"),
    ]
    y = 2.6
    for name, desc in files:
        add_text(slide, Inches(0.85), Inches(y), Inches(4.0), Inches(0.25),
                 name, size=9, bold=True, color=NAVY)
        add_text(slide, Inches(4.9), Inches(y), Inches(4.5), Inches(0.25),
                 desc, size=9, color=DARK_GRAY)
        y += 0.27

    # 技術スタック
    add_round_rect(slide, Inches(0.6), Inches(4.95), Inches(8.8), Inches(0.4), fill=LIGHT_GRAY)
    add_text(slide, Inches(0.85), Inches(5.0), Inches(8.4), Inches(0.3),
             "Python / ChromaDB / sentence-transformers / Gemini API / cc-sdd（仕様駆動開発）",
             size=10, color=DARK_GRAY)


# =====================================================
# スライド構築
# =====================================================
slides_def = [
    slide_title,
    slide_purpose,
    slide_architecture,
    slide_sources,
    slide_rubric,
    slide_question_evolution,
    slide_score_chart,
    slide_v6_questions,
    slide_v6_results,
    slide_why_nlm_works,
    slide_v7_implementation,
    slide_v7_results,
    slide_q2_analysis,
    slide_conclusion,
    slide_lessons,
    slide_future_work,
    slide_references,
]

total_slides = len(slides_def)

for i, build in enumerate(slides_def, start=1):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    build(slide, i, total_slides)
    if i > 1:  # タイトルスライドにはフッター不要
        add_footer(slide, i, total_slides)

output_path = "/Users/nakamurashogo/Desktop/lab/rag-comparison/RAG_vs_NotebookLM_report.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {total_slides}")
